from fastapi.responses import FileResponse
import os
from pathlib import Path
import requests
from datetime import datetime
from pptx import Presentation
from pptx.util  import Pt,Inches
from PIL import Image
from pptx.enum.text import PP_ALIGN  # Import the PP_ALIGN enum
import pandas as pd

def handle_chapter_list_json(chapter_list): 
    print('handle_chapter_list_json start--------------------------Handles a list of chapter dictionaries and formats them.')
    try:
        if not chapter_list or not isinstance(chapter_list, list):
            print("⚠️ Warning: Invalid chapter list provided. Skipping processing.")
            return []
        processed_data = []
        for index, chapter in enumerate(chapter_list):
        # for chapter in chapter_list:
            if index==0:
                print('chapter--------------------------',chapter)
            chapter_desc = chapter.get('chapterDescription', {})
            
            # If chapterDescription is a dict, convert it to a list of strings
            if isinstance(chapter_desc, dict):
                chapter_desc = [f"{k}: {v}" for k, v in chapter_desc.items()]
            else:
                chapter_desc = [str(chapter_desc)]
            
            processed_data.append({
                "chapterSummary": chapter.get("chapterSummary", "No Summary"),
                "chapterDescription": chapter_desc
            })
        return processed_data
    except Exception as e:  
        print(f" Error in handle_has_not_chapters: {e}")
        return None


def handle_dictionary(summary_json):
    print('Handles a single dictionary input for chapter data.--------------------------')
    try:
        processed_data = handle_chapter_list_json(summary_json.get('chapters', []))  # Process chapters list from the dictionary
        return processed_data
    except Exception as e:
        print(f" Error in handle_dictionary: {e}")
        return None
    

def handle_json(summary_json):
    print('handle_json start--------------------------Processes the JSON data and ensures it''s in the correct format.')
    processed_data = []
    try:
        if not summary_json:
            print("⚠️ Warning: No JSON data provided. Skipping processing.")
            return []
        # Check if summary_json is a list of standalone chapters without topics
        if isinstance(summary_json, list) and all("chapterSummary" in item for item in summary_json):
            # Wrap the chapters inside a default topic structure
            default_topic = {
                "topic": "Generated Topic",
                "topic_summary": "This is an auto-generated summary from chapters.",
                "chapters": summary_json }
            processed_data.extend(handle_has_chapters(default_topic))
        else:
            for topic in summary_json:
            # Checking if there is a 'chapters' field and if it contains chapters
                chapters = topic.get("chapters", [])
                
                if not chapters:
                    # No chapters, call the helper for no-chapter processing
                    processed_data.append(handle_has_not_chapters(topic))
                else:
                    # Chapters exist, call the helper for chapters processing
                    processed_data.extend(handle_has_chapters(topic))
               
    except Exception as e:
        print(f" Error in handle_json: {e}")
        return None
    return processed_data

def handle_has_chapters(topic):
    print('handle_has_chapters------------------------Handles topics that have chapters.')
    processed_data = []
    try:
        
        for index, chapter in enumerate(topic.get("chapters", [])):
        # for chapter in topic.get("chapters", []):
            if index==0:
                print('chapter--------------------------', chapter)
            chapter_desc = chapter.get('chapterDescription', [])
            if isinstance(chapter_desc, str):
                chapter_desc = [chapter_desc]
            elif isinstance(chapter_desc, dict):
                chapter_desc = [f"{k}: {v}" for k, v in chapter_desc.items()]
            
            processed_data.append({
                "topic": topic.get("topic", "Unknown Topic"),
                "topic_summary": topic.get("topic_summary", "No topic summary available"),
                "chapterSummary": chapter.get("chapterSummary", "No chapterSummary available"),
                "chapterDescription": chapter_desc  })
    except Exception as e:
        print(f" Error in handle_has_chapters: {e}")
        return None
    return processed_data


def handle_has_not_chapters(topic):
    print('handle_has_not_chapters------------------------ Handles topics that do not have chapters.')
    try:
        chapter_desc = topic.get('chapterDescription', {})
        print('chapter_desc------------------------',chapter_desc)
        
        if isinstance(chapter_desc, dict):
            chapter_desc = [f"{k}: {v}" for k, v in chapter_desc.items()]
        else:
            print('chapter_desc not dict ------------------------')
        return {
            "topic": topic.get("topic", "Unknown Topic"),
            "topic_summary": topic.get("topic_summary", "No topic summary available"),
            "chapterSummary": topic.get("chapterSummary", "No chapterSummary available"),
            "chapterDescription": chapter_desc }
    except Exception as e:
        print(f" Error in handle_has_not_chapters: {e}")
        return None
    
def validate_and_extract_data(data):
    print('start validate_and_extract_data------------------------This function validates the input JSON and ensures it''s in the correct format.')
    try:
        if isinstance(data, dict):
            return [data]  # Convert dictionary to list for uniform processing
        elif isinstance(data, list):
            if not all(isinstance(item, dict) for item in data):
                print("Error: List contains non-dictionary elements.")
                return None
            return data
        else:
            print("Error: Data is neither a dictionary nor a list.")
            return None
    except Exception as e:
        print(f" Error in validate_and_extract_data: {e}")
        return None
        
def handle_ppt(processed_data, ppt_path, file_name, image_filenames):
    print('handle_ppt start--------------------------Creates the PowerPoint presentation using processed data.')
    file_path = os.path.join(ppt_path, file_name + '.pptx')
    image_index = 0    
    print('image_filenames length=',len(image_filenames))
    if image_filenames is None:
        image_filenames = []
    
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  
    try:
        valid_data = validate_and_extract_data(processed_data)
        # print('valid_data---------------------',valid_data)
        
         # -----------------------------------------------------
        # Convert list of dictionaries into a DataFrame
        df = pd.DataFrame(valid_data)

        # Group by 'topic' and aggregate chapters into lists
        grouped_df = df.groupby(["topic", "topic_summary"]).agg(list).reset_index()
        print('grouped_df=',grouped_df)
        # Convert back to a structured list of dictionaries
        grouped_data = grouped_df.to_dict(orient="records")
        print('grouped_data=',grouped_data)
        # -----------------------------------------------------
        create_agenda_slide(prs, valid_data)  # Add agenda slide first
        
        if not valid_data:
            raise ValueError("Invalid JSON data")
        for index, data in enumerate(valid_data): 
            print(f'Processing data[{index}] =', data)

            slide = prs.slides.add_slide(slide_layout)
            title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
            title_frame = title_shape.text_frame
            title_frame.word_wrap = True  # Enable word wrapping
            
            topic = data.get('topic', 'Untitled')
            topic_summary = data.get('topic_summary', '')
            title_frame.text = f"{topic} - {topic_summary}" if topic_summary else f"{topic}"
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.size = Pt(22)
            text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5), Inches(5))
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            
            
            # Set margins to avoid text getting cut off
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            text_frame.margin_top = Inches(0.1)
            text_frame.margin_bottom = Inches(0.1)

            for paragraph in text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.JUSTIFY

            p1 = text_frame.add_paragraph()    
            
            try:
                chapter_summary = data.get('chapterSummary', 'No Chapter Summary')
                p1.text = chapter_summary
            except Exception as e:
                print(f"Error in handle_has_not_chapters: {e}")
            p1.font.bold = True
            p1.font.size = Pt(18)
            p1.alignment = PP_ALIGN.LEFT
            
            chapter_description = data.get('chapterDescription', [])
            if isinstance(chapter_description, list):
                for desc in chapter_description:
                    p2 = text_frame.add_paragraph()
                    p2.text = desc.strip()
                    p2.font.size = Pt(16)
                    p2.alignment = PP_ALIGN.LEFT
            else:
                print(f"Warning: chapterDescription is not a list in {data}")
                
            if image_index < len(image_filenames):
                
                image_path = image_filenames[image_index]
                if image_path and os.path.exists(image_path):
                    
                    left, top, width, height = Inches(5.5), Inches(1.5), Inches(4), Inches(5)
                    img = Image.open(image_path)
                    img.thumbnail((1000, int((1000 / img.width) * img.height)), Image.Resampling.LANCZOS)
                    temp_image_path = 'temp_resized_image.jpg'
                    img.save(temp_image_path)
                    slide.shapes.add_picture(temp_image_path, left, top, width, height)
                    os.remove(temp_image_path)
                    print('image added=',image_index)
                image_index += 1
        print(image_index)
        add_remaining_images(image_filenames,image_index,prs,slide_layout)
       
        prs.save(file_path)  
        # headers = { "Content-Disposition": f'attachment; filename="{file_name}.pptx"'}
        # response = FileResponse(
        #     path=file_path,  
        #     media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        #     headers=headers )
        # print(f"Presentation created successfully at {datetime.now().time()}")
        # return response
        return file_path
    except Exception as e:
        print(f"Error in handle_ppt: {e}")
        return None
def add_remaining_images(image_filenames,image_index,prs,slide_layout):
    print('add_remaining_images start')
    ramaining_images=image_filenames[image_index:]
    i=0
    while i< len(ramaining_images):
        print('add_remaining_images=',i)
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
        title_frame = title_shape.text_frame
        title_frame.text = "Additional images"
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.size = Pt(22)
        
        position=[(Inches(0.5),Inches(1.5)),(Inches(5.5),Inches(1.5))]
        for j in range(2):
            if i<len(ramaining_images):
                image_path=ramaining_images[i]
                if image_path and os.path.exists(image_path):
                    left, top=position[j]
                    width, height=Inches(4), Inches(5)
                    img = Image.open(image_path)
                    img.thumbnail((1000, int((1000 / img.width) * img.height)), Image.Resampling.LANCZOS)
                    temp_image_path = 'temp_resized_image_{i}.jpg'
                    img.save(temp_image_path)
                    
                    slide.shapes.add_picture(temp_image_path, left, top, width, height)
                    os.remove(temp_image_path)
                    print(f'image_index added{i+1}/{len(ramaining_images)} at position{j}')
                i += 1    
    return prs  


def create_agenda_slide(prs, processed_data):
    
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
        
    for data in processed_data:
        topic = data.get("topic", "Untitled")
        chapters = [data["chapterSummary"]] if "chapterSummary" in data else []
       
        title.text = f"Agenda: {topic}"
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.size = Pt(32)

        for chapter in chapters:
            p = content.text_frame.add_paragraph()
            p.text = f"{chapter}"
            p.font.size = Pt(20)
            p.alignment = PP_ALIGN.LEFT
    
    print("Agenda slide created successfully")

def convert_json_to_ppt(request, summary_json, ppt_path, file_name, image_filenames):
    # print('convert_json_to_ppt start----------------------------------')
    try:
        # Check if summary_json is a list (multiple topics) or a single dictionary
        if isinstance(summary_json, list):
            print('found list----------------------------------')
            processed_data = handle_json(summary_json) 
            # print('handle_json end----------------------------------')
        elif isinstance(summary_json, dict):
            print('found dicttionary----------------------------------')
            processed_data = handle_dictionary(summary_json)
            # print('handle_dictionary end----------------------------------')
        else:
            print("⚠️ Invalid JSON format")
            return None
        
        if not processed_data:
            return None
        return handle_ppt(processed_data, ppt_path, file_name, image_filenames)
    except RuntimeError as e:
        raise RuntimeError(f"Error in json formate response:{e}")
    except Exception as e:
        print(f" Error in convert_json_to_ppt: {e}")
        return None
